#!/usr/bin/env python3
"""
Create multiple test PowerPoint presentations for batch testing
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import sys

def create_test_presentation(filename, title, slides_data):
    """Create a test PowerPoint presentation with specified content."""
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    subtitle_placeholder.text = "Test presentation for batch conversion"
    
    # Content slides
    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title_placeholder.text = slide_data["title"]
        
        # Add bullet points
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, bullet in enumerate(slide_data["bullets"]):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
    
    # Save presentation
    prs.save(filename)
    print(f"Created: {filename}")

def main():
    # Create test presentations folder
    test_folder = Path("test_presentations")
    test_folder.mkdir(exist_ok=True)
    
    # Test presentation 1: Python Basics
    create_test_presentation(
        test_folder / "python_basics.pptx",
        "Python Programming Basics",
        [
            {
                "title": "Variables and Data Types",
                "bullets": [
                    ">>> name = 'John'",
                    ">>> age = 25",
                    ">>> height = 5.9",
                    ">>> is_student = True"
                ]
            },
            {
                "title": "Lists and Loops",
                "bullets": [
                    ">>> numbers = [1, 2, 3, 4, 5]",
                    ">>> for num in numbers:",
                    "...     print(num * 2)",
                    "2",
                    "4",
                    "6",
                    "8",
                    "10"
                ]
            }
        ]
    )
    
    # Test presentation 2: Data Structures
    create_test_presentation(
        test_folder / "data_structures.pptx",
        "Data Structures in Python",
        [
            {
                "title": "Dictionaries",
                "bullets": [
                    ">>> student = {'name': 'Alice', 'grade': 'A'}",
                    ">>> print(student['name'])",
                    "Alice",
                    ">>> student['age'] = 20"
                ]
            },
            {
                "title": "Sets",
                "bullets": [
                    ">>> fruits = {'apple', 'banana', 'orange'}",
                    ">>> fruits.add('grape')",
                    ">>> print(len(fruits))",
                    "4"
                ]
            }
        ]
    )
    
    # Test presentation 3: Functions
    create_test_presentation(
        test_folder / "functions.pptx",
        "Python Functions",
        [
            {
                "title": "Defining Functions",
                "bullets": [
                    ">>> def greet(name):",
                    "...     return f'Hello, {name}!'",
                    ">>> greet('World')",
                    "'Hello, World!'"
                ]
            },
            {
                "title": "Lambda Functions",
                "bullets": [
                    ">>> square = lambda x: x ** 2",
                    ">>> square(5)",
                    "25",
                    ">>> numbers = [1, 2, 3, 4]",
                    ">>> list(map(square, numbers))",
                    "[1, 4, 9, 16]"
                ]
            }
        ]
    )
    
    print(f"\n✅ Created 3 test PowerPoint presentations in {test_folder}")
    print("You can now test the batch converter with:")
    print(f"python batch_convert.py {test_folder}")

if __name__ == "__main__":
    main()
