#!/usr/bin/env python3
"""
PowerPoint to Markdown Converter

This script converts a PowerPoint presentation (.pptx) to a Markdown (.md) file.
It extracts text from slides, preserves formatting for tables and bullet points,
and organizes content by slide number.

Usage:
    python pptx_to_md.py <input_file.pptx> [output_file.md]

Requirements:
    - python-pptx library (install via: pip install python-pptx)
"""

import argparse
import sys
import os
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def format_python_code_blocks(text):
    """
    Convert Python interpreter examples to proper code blocks.
    Handles bullet points with >>> commands followed by output bullet points.
    
    Args:
        text: Text content that may contain Python interpreter examples
        
    Returns:
        str: Formatted text with Python code blocks
    """
    if not text or '>>>' not in text:
        return text
    
    lines = text.split('\n')
    result_lines = []
    current_code_block = []
    in_code_section = False
    i = 0
    
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        
        # Check if this is a bullet point with Python interpreter code
        if stripped.startswith('- >>>'):
            # Start a new code block if not already in one
            if not in_code_section:
                in_code_section = True
                current_code_block = []
            
            # Extract the Python command (remove "- >>> ")
            command = stripped[5:].strip()  # Remove "- >>> "
            if command:
                current_code_block.append(command)
            
            # Look ahead for the output (next line should be the result)
            if i + 1 < len(lines):
                next_line = lines[i + 1].strip()
                if next_line.startswith('- ') and not next_line.startswith('- >>>'):
                    # This is likely the output
                    output = next_line[2:].strip()  # Remove "- "
                    if output and len(output) < 100:  # Reasonable output length
                        current_code_block.append(f"# Output: {output}")
                    i += 1  # Skip the output line
            
        elif stripped.startswith('- ') and not stripped.startswith('- >>>') and in_code_section:
            # This might be additional output or we're moving to regular bullet points
            output = stripped[2:].strip()
            
            # Simple heuristic: if it looks like output (short, simple), include it
            if (len(output) < 50 and 
                (output.replace('.', '').replace('-', '').isdigit() or
                 output in ['True', 'False'] or
                 output.startswith("'") or output.startswith('"') or
                 'Error' in output or 'Traceback' in output)):
                current_code_block.append(f"# Output: {output}")
            else:
                # End the code block and start regular content
                if current_code_block:
                    result_lines.append("```python")
                    result_lines.extend(current_code_block)
                    result_lines.append("```")
                    result_lines.append("")  # Add blank line
                    current_code_block = []
                in_code_section = False
                result_lines.append(line)
        
        else:
            # Regular line - end any current code block
            if in_code_section and current_code_block:
                result_lines.append("```python")
                result_lines.extend(current_code_block)
                result_lines.append("```")
                result_lines.append("")  # Add blank line
                current_code_block = []
            in_code_section = False
            
            # Add the regular line
            if stripped or line.startswith(('- ', '  -')):
                result_lines.append(line)
        
        i += 1
    
    # Handle any remaining code block
    if in_code_section and current_code_block:
        result_lines.append("```python")
        result_lines.extend(current_code_block)
        result_lines.append("```")
    
    return '\n'.join(result_lines)


def extract_text_from_shape(shape):
    """
    Extract text from a shape, handling different shape types.
    
    Args:
        shape: A shape object from python-pptx
        
    Returns:
        str: Extracted text content
    """
    text_content = ""
    
    if hasattr(shape, "text_frame") and shape.text_frame:
        paragraphs = []
        for paragraph in shape.text_frame.paragraphs:
            para_text = paragraph.text.strip()
            if para_text:
                # Check if this is a bulleted or indented paragraph
                if paragraph.level > 0:
                    indent = "  " * paragraph.level
                    paragraphs.append(f"{indent}- {para_text}")
                elif len(shape.text_frame.paragraphs) > 1 and paragraph != shape.text_frame.paragraphs[0]:
                    # This is a subsequent paragraph in a multi-paragraph text frame
                    paragraphs.append(f"- {para_text}")
                else:
                    # This is a title or single paragraph
                    paragraphs.append(para_text)
        
        if paragraphs:
            text_content = "\n".join(paragraphs)
    elif hasattr(shape, "text") and shape.text.strip():
        text_content = shape.text.strip()
    
    # Format Python code blocks if the text contains interpreter examples
    if text_content and '>>>' in text_content:
        text_content = format_python_code_blocks(text_content)
    
    return text_content


def extract_table_content(table):
    """
    Extract content from a table and format it as Markdown.
    
    Args:
        table: A table object from python-pptx
        
    Returns:
        str: Markdown formatted table
    """
    if not table:
        return ""
    
    markdown_table = []
    
    # Extract table data
    table_data = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            cell_text = cell.text.strip() if cell.text else ""
            # Replace newlines in cells with spaces for better markdown formatting
            cell_text = cell_text.replace('\n', ' ').replace('\r', ' ')
            row_data.append(cell_text)
        table_data.append(row_data)
    
    if not table_data:
        return ""
    
    # Format as Markdown table
    if len(table_data) > 0:
        # Header row
        header = "| " + " | ".join(table_data[0]) + " |"
        markdown_table.append(header)
        
        # Separator row
        separator = "| " + " | ".join(["---"] * len(table_data[0])) + " |"
        markdown_table.append(separator)
        
        # Data rows
        for row in table_data[1:]:
            row_str = "| " + " | ".join(row) + " |"
            markdown_table.append(row_str)
    
    return "\n".join(markdown_table)


def generate_table_of_contents(presentation):
    """
    Generate a table of contents with links to each slide.
    
    Args:
        presentation: A Presentation object from python-pptx
        
    Returns:
        str: Markdown table of contents
    """
    num_slides = len(presentation.slides)
    
    # For long presentations, make TOC collapsible
    if num_slides > 20:
        toc_lines = ["<details>", "<summary><strong>📋 Table of Contents</strong> (Click to expand)</summary>", "", "## Table of Contents", ""]
    else:
        toc_lines = ["## Table of Contents", ""]
    
    for i, slide in enumerate(presentation.slides, 1):
        # Try to extract a meaningful title for each slide
        slide_title = f"Slide {i}"
        
        # Look for the best title in the slide shapes
        potential_titles = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                if text and len(text) < 80:  # Reasonable title length
                    clean_title = text.replace('\n', ' ').replace('\r', ' ').strip()
                    # Filter out unwanted content
                    if (not clean_title.startswith(('- ', '• ', '>>>', 'Traceback', 'File', 'http', 'www')) and 
                        not any(word in clean_title.lower() for word in ['output:', 'error:', 'exception', 'class \'']) and
                        len(clean_title.split()) <= 6 and  # Max 6 words for title
                        not clean_title.endswith(('.com', '.org', '.net'))):  # Filter out URLs
                        potential_titles.append((clean_title, len(clean_title.split())))
                        
            elif hasattr(shape, "text_frame") and shape.text_frame:
                # Check first paragraph for title
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        text = paragraph.text.strip()
                        if text and len(text) < 80:
                            clean_title = text.replace('\n', ' ').replace('\r', ' ').strip()
                            if (not clean_title.startswith(('- ', '• ', '>>>', 'Traceback', 'File', 'http', 'www')) and 
                                not any(word in clean_title.lower() for word in ['output:', 'error:', 'exception', 'class \'']) and
                                len(clean_title.split()) <= 6 and
                                not clean_title.endswith(('.com', '.org', '.net'))):
                                potential_titles.append((clean_title, len(clean_title.split())))
                        break  # Only check first paragraph
        
        # Choose the best title
        if potential_titles:
            # Sort by word count (prefer shorter titles) and take the first one
            potential_titles.sort(key=lambda x: x[1])
            best_title = potential_titles[0][0]
            
            # Special handling for common generic titles
            generic_titles = ['introduction to computing using python', 'slide', 'exercise', 'example']
            if best_title.lower() in generic_titles and len(potential_titles) > 1:
                # Look for more specific content
                for title, word_count in potential_titles[1:]:
                    if title.lower() not in generic_titles:
                        best_title = title
                        break
            
            slide_title = f"Slide {i}: {best_title}"
        
        # Create markdown link to slide
        slide_anchor = f"slide-{i}"
        toc_lines.append(f"{i}. [{slide_title}](#{slide_anchor})")
    
    # Close collapsible section if used
    if num_slides > 20:
        toc_lines.append("")
        toc_lines.append("</details>")
    
    toc_lines.append("")  # Empty line after TOC
    return "\n".join(toc_lines)


def process_slide(slide, slide_number):
    """
    Process a single slide and extract all text content.
    
    Args:
        slide: A slide object from python-pptx
        slide_number: The slide number (1-indexed)
        
    Returns:
        str: Markdown formatted content for the slide
    """
    # Create slide header with anchor for TOC linking
    slide_content = [f'<a id="slide-{slide_number}"></a>']
    slide_content.append(f"## Slide {slide_number}")
    
    # Track if we found any content
    has_content = False
    
    for shape in slide.shapes:
        content = ""
        
        # Handle different shape types
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            # Extract table content
            content = extract_table_content(shape.table)
            if content:
                slide_content.append("\n" + content + "\n")
                has_content = True
                
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or \
             shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or \
             shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            # Extract text content
            content = extract_text_from_shape(shape)
            if content:
                slide_content.append(content)
                has_content = True
                
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Handle grouped shapes
            for grouped_shape in shape.shapes:
                grouped_content = extract_text_from_shape(grouped_shape)
                if grouped_content:
                    slide_content.append(grouped_content)
                    has_content = True
        
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Note the presence of images
            slide_content.append("*[Image present]*")
            has_content = True
    
    # If no content was found, add a note
    if not has_content:
        slide_content.append("*[No text content found]*")
    
    return "\n\n".join(slide_content)


def convert_pptx_to_markdown(pptx_path, output_path=None):
    """
    Convert a PowerPoint presentation to Markdown format.
    
    Args:
        pptx_path (str): Path to the input .pptx file
        output_path (str, optional): Path for the output .md file
        
    Returns:
        str: Path to the created markdown file
    """
    # Validate input file
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"Input file not found: {pptx_path}")
    
    if not pptx_path.lower().endswith('.pptx'):
        raise ValueError("Input file must be a .pptx file")
    
    # Determine output path
    if output_path is None:
        base_name = Path(pptx_path).stem
        output_path = f"{base_name}.md"
    
    try:
        # Load the presentation
        print(f"Loading presentation: {pptx_path}")
        presentation = Presentation(pptx_path)
        
        # Start building the markdown content
        markdown_content = []
        
        # Add header
        presentation_name = Path(pptx_path).stem
        markdown_content.append(f"# {presentation_name}")
        markdown_content.append(f"Converted from PowerPoint presentation: `{Path(pptx_path).name}`")
        markdown_content.append(f"Total slides: {len(presentation.slides)}")
        markdown_content.append("")  # Add blank line
        
        # Generate and add table of contents
        print("Generating table of contents...")
        toc = generate_table_of_contents(presentation)
        markdown_content.append(toc)
        
        # Process each slide
        print(f"Processing {len(presentation.slides)} slides...")
        for i, slide in enumerate(presentation.slides, 1):
            print(f"Processing slide {i}...")
            slide_content = process_slide(slide, i)
            markdown_content.append(slide_content)
        
        # Write to output file
        final_content = "\n\n---\n\n".join(markdown_content)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(final_content)
        
        print(f"Conversion completed! Output saved to: {output_path}")
        return output_path
        
    except Exception as e:
        print(f"Error processing presentation: {str(e)}")
        raise


def main():
    """Main function to handle command line arguments and run the conversion."""
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint presentation to Markdown format",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
    python pptx_to_md.py presentation.pptx
    python pptx_to_md.py presentation.pptx output.md
        """
    )
    
    parser.add_argument(
        'input_file',
        help='Path to the input PowerPoint file (.pptx)'
    )
    
    parser.add_argument(
        'output_file',
        nargs='?',
        help='Path to the output Markdown file (.md). If not specified, uses input filename with .md extension'
    )
    
    parser.add_argument(
        '--version',
        action='version',
        version='PowerPoint to Markdown Converter 2.0'
    )
    
    args = parser.parse_args()
    
    try:
        output_file = convert_pptx_to_markdown(args.input_file, args.output_file)
        print(f"Successfully converted '{args.input_file}' to '{output_file}'")
        
    except FileNotFoundError as e:
        print(f"Error: {e}")
        sys.exit(1)
        
    except ValueError as e:
        print(f"Error: {e}")
        sys.exit(1)
        
    except Exception as e:
        print(f"Unexpected error: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()