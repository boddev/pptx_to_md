#!/usr/bin/env python3
"""
Test script for the PowerPoint to Markdown converter.
Creates a sample PowerPoint presentation for testing purposes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os


def create_test_presentation():
    """Create a sample PowerPoint presentation for testing."""
    
    # Create a presentation object
    prs = Presentation()
    
    # Slide 1: Title slide with bullet points
    slide1_layout = prs.slide_layouts[1]  # Title and Content layout
    slide1 = prs.slides.add_slide(slide1_layout)
    
    title1 = slide1.shapes.title
    title1.text = "Welcome to Our Product Demo"
    
    content1 = slide1.placeholders[1]
    tf1 = content1.text_frame
    tf1.text = "Key Features:"
    
    p = tf1.add_paragraph()
    p.text = "User-friendly interface"
    p.level = 1
    
    p = tf1.add_paragraph()
    p.text = "Advanced analytics"
    p.level = 1
    
    p = tf1.add_paragraph()
    p.text = "Cloud integration"
    p.level = 1
    
    p = tf1.add_paragraph()
    p.text = "Real-time collaboration"
    p.level = 2
    
    p = tf1.add_paragraph()
    p.text = "Multi-platform support"
    p.level = 2
    
    # Slide 2: Table slide
    slide2_layout = prs.slide_layouts[6]  # Blank layout
    slide2 = prs.slides.add_slide(slide2_layout)
    
    # Add title
    title2 = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title2.text_frame.text = "Pricing Comparison"
    title2.text_frame.paragraphs[0].font.size = Pt(24)
    title2.text_frame.paragraphs[0].font.bold = True
    
    # Add table
    rows, cols = 4, 3
    table = slide2.shapes.add_table(rows, cols, Inches(2), Inches(2), Inches(6), Inches(3))
    
    # Header row
    table.table.cell(0, 0).text = "Plan"
    table.table.cell(0, 1).text = "Features"
    table.table.cell(0, 2).text = "Price"
    
    # Data rows
    table.table.cell(1, 0).text = "Basic"
    table.table.cell(1, 1).text = "Core features"
    table.table.cell(1, 2).text = "$9/month"
    
    table.table.cell(2, 0).text = "Pro"
    table.table.cell(2, 1).text = "Advanced analytics"
    table.table.cell(2, 2).text = "$19/month"
    
    table.table.cell(3, 0).text = "Enterprise"
    table.table.cell(3, 1).text = "Full suite + support"
    table.table.cell(3, 2).text = "$49/month"
    
    # Slide 3: Text with multiple text boxes
    slide3_layout = prs.slide_layouts[6]  # Blank layout
    slide3 = prs.slides.add_slide(slide3_layout)
    
    # Add title
    title3 = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title3.text_frame.text = "Implementation Timeline"
    title3.text_frame.paragraphs[0].font.size = Pt(24)
    title3.text_frame.paragraphs[0].font.bold = True
    
    # Add multiple text boxes
    textbox1 = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(2))
    tf = textbox1.text_frame
    tf.text = "Phase 1: Planning"
    p = tf.add_paragraph()
    p.text = "Requirements gathering"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Resource allocation"
    p.level = 1
    
    textbox2 = slide3.shapes.add_textbox(Inches(5), Inches(2), Inches(3), Inches(2))
    tf2 = textbox2.text_frame
    tf2.text = "Phase 2: Development"
    p = tf2.add_paragraph()
    p.text = "System setup"
    p.level = 1
    p = tf2.add_paragraph()
    p.text = "Feature implementation"
    p.level = 1
    
    textbox3 = slide3.shapes.add_textbox(Inches(3), Inches(4.5), Inches(3), Inches(1.5))
    tf3 = textbox3.text_frame
    tf3.text = "Phase 3: Testing & Deployment"
    
    # Save the presentation
    output_file = "test_presentation.pptx"
    prs.save(output_file)
    print(f"Test presentation created: {output_file}")
    return output_file


if __name__ == "__main__":
    print("Creating test PowerPoint presentation...")
    test_file = create_test_presentation()
    
    print(f"\nTest file created: {test_file}")
    print("You can now test the converter with:")
    print(f"python pptx_to_md.py {test_file}")
